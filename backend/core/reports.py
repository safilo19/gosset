"""export_report logic: renders prior structuredContent results to Markdown, xlsx, Word or PDF.

Plain Python — no MCP or web-framework code.

Word uses python-docx and PDF uses reportlab: both are pure-Python wheels. That is deliberate —
the weasyprint/GTK family needs system libraries that don't install cleanly on Windows, and this
app runs locally on Windows as often as it runs on Render.
"""

from __future__ import annotations

import base64
import io
import json
import os
import re
import shutil
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

from PIL import Image as PILImage
from docx import Document
from docx.enum.section import WD_ORIENT, WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_TAB_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from backend import version
from backend.core import charts as charts_core
from backend.report_engine import builder as report_builder
from backend.report_engine import components as report_components
from backend.report_engine import theme as report_theme
from backend.report_engine import verdict as verdict_engine

# backend/core/reports.py -> backend/core -> backend -> project root.
#
# Under PyInstaller this resolves to the unpacked bundle directory, which is exactly right for
# READING bundled assets (_MARK_PNG below): the .spec ships frontend/ at the bundle root, so every
# `PROJECT_ROOT / "frontend" / ...` path in this package keeps working unchanged when frozen.
PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent

# WRITING is the other half, and it cannot use PROJECT_ROOT when frozen: that would put generated
# reports inside the installed application directory. GOSSET_OUTPUT_DIR is how the desktop shell
# points this at a per-user location; a source checkout keeps ./output as before.
OUTPUT_DIR = Path(os.environ.get("GOSSET_OUTPUT_DIR") or (PROJECT_ROOT / "output"))

# Branding on exports stays to a header/footer credit line, deliberately: a report belongs to
# whoever ran the analysis, so the app signs it rather than stamping it. No watermark, no cover
# page.
APP_NAME = "Gosset"
APP_VERSION = version.VERSION  # generated from desktop/package.json; About shows the same string
CREDIT_LINE = f"Generated with {APP_NAME}"
# The accent mark, already rasterised for the favicons. Reused here so a document header needs no
# SVG rasteriser; None when it has not been generated (frontend/brand/make_favicons.py writes it).
_MARK_PNG = PROJECT_ROOT / "frontend" / "brand" / "favicon-48.png"


def _mark_png() -> Path | None:
    return _MARK_PNG if _MARK_PNG.is_file() else None


_TABLE_KEYS_PRIORITY = [
    "coefficients",
    "forecast",
    "segments",
    "strongest_pairs",
    "stats",
    "groups",
    "datasets",
    "preview",
    "feature_importances",
    "results",
]

# Fields that are never useful in a rendered key/value table: raw image bytes (the chart is already
# embedded as an actual image via chart_path) and anything else absurdly long for a table cell.
_SCALAR_EXCLUDE_KEYS = {"image_base64", "chart_path"}
_SCALAR_MAX_LEN = 300

# Plain-language narrative fields (already present across the analysis outputs) get pulled out
# of the key/value table and rendered as ordinary paragraph text instead — a wall of prose
# crammed into a two-column table cell reads worse than the same sentence set as a paragraph.
_NARRATIVE_KEYS = ("conclusion", "interpretation", "summary", "method_reason")

# Design tokens for the Word and Excel exports, as hex strings because that is what python-docx and
# openpyxl want. The PDF's palette is not here — it lives in report_engine/theme.py, and these two
# read FROM it so a retune of the report engine's colours reaches Word and Excel as well.
_INK = report_theme.INK.hexval()[2:]
_MUTED = report_theme.MUTED.hexval()[2:]
_ACCENT = report_theme.ACCENT.hexval()[2:]
_SUCCESS = report_theme.SUCCESS.hexval()[2:]
_DANGER = report_theme.DANGER.hexval()[2:]
_BORDER = report_theme.BORDER.hexval()[2:]
_HAIRLINE = report_theme.HAIRLINE.hexval()[2:]
_HEADER_BG = report_theme.HEADER_BG.hexval()[2:]
_ROW_ALT_BG = report_theme.ROW_ALT_BG.hexval()[2:]
_INK_80 = report_theme.INK_80.hexval()[2:]

DEFAULT_DECIMALS = 3

# A report is a document, not a data dump. Past TABLE_ROW_LIMIT rows a table stops being readable and
# starts paginating across pages — one 96-row table used to take three pages of a report and bury
# every finding after it. Such a table exports as its first TABLE_ROWS_SHOWN rows plus a mono note
# saying how many were left out. The Report pane's "include full table" option sets
# ReportSection.full_tables to opt one section out, which is the only way to get the long form.
#
# The EXCEL export is deliberately exempt: a spreadsheet is where someone goes FOR the rows, and
# truncating there would be taking away the one format that should hold everything.
TABLE_ROW_LIMIT = 25
TABLE_ROWS_SHOWN = 20


def _document_rows(rows: list[dict[str, Any]] | None, *, allow_full: bool = False) -> tuple[list[dict[str, Any]], str]:
    """(rows to render, mono note or "") for a table going into a document."""
    rows = rows or []
    if allow_full or len(rows) <= TABLE_ROW_LIMIT:
        return rows, ""
    hidden = len(rows) - TABLE_ROWS_SHOWN
    return rows[:TABLE_ROWS_SHOWN], f"… {hidden} more rows — full table available in {APP_NAME}"


@dataclass
class ReportSection:
    title: str
    data: dict[str, Any]
    chart_path: str | None = None
    # A PNG data URL (or bare base64) captured from the chart the user actually saw on screen.
    # Any chart type can supply one, which is how plugin and plotly charts reach the exports
    # without a server-side renderer per type.
    chart_image_base64: str | None = None
    # False for a section that IS one output block: a table block must not also acquire a chart
    # rendered from its own rows, which is what section_chart's last-resort branch would do.
    allow_generated_chart: bool = True
    # True when the person explicitly staged this block as a FULL table ("include full table" on the
    # block's Send-to-Report action). Otherwise a long table is truncated — see TABLE_ROW_LIMIT.
    full_tables: bool = False
    # Which analysis produced this, so the verdict badge knows whether a small p-value is a finding
    # (a t-test) or a failure (a normality test). See report_engine/verdict.py.
    analysis_id: str = ""
    # Set for a Report-pane note: commentary from the user, not a result. Renders as prose with no
    # heading, badge or table, so it reads as the author speaking between the findings.
    note: str = ""
    columns: str = ""  # the input columns, for the card's metadata line
    timestamp: str = ""  # when the analysis ran, already formatted
    # filled in on first use, so every format in one run embeds the same generated PNG
    generated_chart: Path | None = field(default=None, repr=False)


@dataclass
class ReportMeta:
    """Everything the title block needs. `decimals` comes from the app's File > Options setting so
    the numbers in an exported report match the numbers on screen."""

    dataset_id: str
    source: str = ""
    row_count: int | None = None
    column_count: int | None = None
    decimals: int = DEFAULT_DECIMALS

    @property
    def dataset_label(self) -> str:
        return self.source or self.dataset_id

    @property
    def title(self) -> str:
        return f"{self.dataset_label} — Analysis Report"

    @property
    def shape_line(self) -> str:
        bits = [f"Dataset {self.dataset_id}"]
        if self.row_count is not None and self.column_count is not None:
            bits.append(f"{self.row_count} rows × {self.column_count} columns")
        bits.append(datetime.now().strftime("Generated %d %B %Y at %H:%M"))
        return "  ·  ".join(bits)


def coerce_meta(meta: ReportMeta | str) -> ReportMeta:
    """Callers used to pass a bare dataset_id; accept either."""
    return meta if isinstance(meta, ReportMeta) else ReportMeta(dataset_id=str(meta))


def default_stem(dataset_id: str) -> str:
    return f"report_{dataset_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"


def summarize_export(dataset_id: str, section_titles: list[str], paths: list[Path]) -> str:
    return f"Exported {len(section_titles)} section(s) for dataset '{dataset_id}' to: {', '.join(str(p) for p in paths)}."


def _slugify(text: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", text.strip().lower()).strip("_")
    return slug or "section"


def _sheet_name(title: str, index: int, used_names: set[str]) -> str:
    cleaned = re.sub(r"[:\\/?*\[\]]", "-", title).strip()[:28] or f"Sheet{index}"
    name = cleaned
    suffix = 1
    while name in used_names:
        name = f"{cleaned[:25]}_{suffix}"
        suffix += 1
    return name


# ---------------------------------------------------------------------------
# shared: value formatting and table extraction
# ---------------------------------------------------------------------------


def _is_number(value: Any) -> bool:
    return isinstance(value, (int, float)) and not isinstance(value, bool)


def format_number(value: float, decimals: int) -> str:
    """One rule for every format, so a value reads the same in the PDF, the Word file and Excel.
    Below 10^-decimals a fixed rendering would collapse to 0.000 and throw the answer away
    (p-values especially), so those switch to exponential."""
    if isinstance(value, int):
        return str(value)
    if float(value).is_integer():
        return str(int(value))
    if value != 0 and abs(value) < 10 ** -decimals:
        return f"{value:.2e}"
    return f"{value:.{decimals}f}"


def cell_text(value: Any, decimals: int = DEFAULT_DECIMALS) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "yes" if value else "no"
    if _is_number(value):
        return format_number(value, decimals)
    if isinstance(value, (list, dict)):
        return json.dumps(value)
    return str(value)


def _flatten_row(row: dict[str, Any]) -> dict[str, Any]:
    """A row whose value is itself a dict (e.g. segmentation's mean_values) renders as an opaque
    JSON blob in a table cell; flatten it so each nested value gets its own column instead."""
    flat: dict[str, Any] = {}
    for k, v in row.items():
        if isinstance(v, dict):
            flat.update(v)
        else:
            flat[k] = v
    return flat


def _extract_table(data: dict[str, Any]) -> tuple[str | None, list[dict[str, Any]], dict[str, Any]]:
    """Pull the first known list-of-dicts (or matrix-like dict) out of `data` as the main table;
    everything else scalar goes in a separate key/value table."""
    table_key: str | None = None
    rows: list[dict[str, Any]] = []

    for key in _TABLE_KEYS_PRIORITY:
        value = data.get(key)
        if isinstance(value, list) and value and all(isinstance(item, dict) for item in value):
            table_key = key
            rows = [_flatten_row(item) for item in value]
            break

    if table_key is None:
        for key in ("matrix", "contingency_table"):
            value = data.get(key)
            if isinstance(value, dict) and value:
                table_key = key
                rows = [{"": row_key, **row_val} for row_key, row_val in value.items()]
                break

    # The length cap keeps runaway strings out of a two-column table cell — but the narrative
    # fields are exactly the long ones, and they are rendered as paragraphs rather than cells,
    # so they must not be capped away. (They were, which silently dropped every regression's
    # interpretation from the report.)
    scalars = {}
    for k, v in data.items():
        if k == table_key or k in _SCALAR_EXCLUDE_KEYS or isinstance(v, (list, dict)):
            continue
        if k in _NARRATIVE_KEYS or len(str(v)) <= _SCALAR_MAX_LEN:
            scalars[k] = v
    return table_key, rows, scalars


def _extract_named_tables(data: dict[str, Any]) -> list[tuple[str, list[dict[str, Any]]]]:
    """[(title, rows)] for a result that ships several output blocks rather than one table.
    Stat > Basic Statistics does: Minitab prints Descriptive Statistics, Estimation and Test as
    separate tables, and each keeps its own heading here."""
    named: list[tuple[str, list[dict[str, Any]]]] = []
    for entry in data.get("tables") or []:
        if not isinstance(entry, dict):
            continue
        rows = entry.get("rows")
        if isinstance(rows, list) and rows and all(isinstance(row, dict) for row in rows):
            named.append((str(entry.get("title") or ""), [_flatten_row(row) for row in rows]))
    return named


def _caption_key(text: str) -> str:
    """A caption fragment reduced to what it SAYS, for comparing two of them.

    Whitespace collapsed, SEPARATOR punctuation flattened to a space, trailing punctuation dropped,
    case folded. The raw comparison this replaces was a plain substring test, which missed both of the
    pairs that actually turned up in reports:

    * an interval plot emits `conclusion` "Interval Plot of y by g." and `summary` "Interval Plot of
      y." — one is not a substring of the other only because of the full stop;
    * a Calc result's title is "Random Data: 120 rows from Normal(...)" while its `summary` says
      "Random Data — 120 rows from Normal(...)" — the same sentence with a different separator.

    Flattening ':' '—' '·' and friends is what makes the second pair compare equal, so the caption
    stops repeating the heading in slightly different punctuation.
    """
    flattened = re.sub(r"[\s]*[—–:;·|-]+[\s]*", " ", re.sub(r"\s+", " ", text))
    return re.sub(r"[\s.,]+$", "", flattened.strip()).casefold()


def _split_narrative(scalars: dict[str, Any], title: str = "") -> tuple[str, dict[str, Any]]:
    """Join the plain-language fields into ONE caption paragraph.

    The fields overlap by design — segmentation's `interpretation` already quotes its `method_reason` —
    so a fragment that says what an earlier one already said is dropped rather than printed twice, and
    a fragment that is merely the section's own TITLE restated is dropped too: a caption's job is to
    say something the heading above it does not. Both tests run on `_caption_key`, so punctuation and
    casing cannot smuggle a repeat through.
    """
    parts: list[str] = []
    keys: list[str] = []
    title_key = _caption_key(title)
    # "Interval Plot: y versus g" -> also compare against "y versus g", so a title that merely names
    # the procedure and its columns still matches a caption that names the same thing.
    title_tail = _caption_key(title.split(":", 1)[1]) if ":" in title else ""
    for key in _NARRATIVE_KEYS:
        value = scalars.get(key)
        if not value:
            continue
        text = str(value).strip()
        if not text:
            continue
        folded = _caption_key(text)
        if not folded or folded == title_key or (title_tail and folded == title_tail):
            continue
        if any(folded in existing or existing in folded for existing in keys):
            continue
        parts.append(text)
        keys.append(folded)
    rest = {k: v for k, v in scalars.items() if k not in _NARRATIVE_KEYS}
    return " ".join(parts), rest


def _is_note_section(section: ReportSection) -> bool:
    """True for a Report-pane note rather than an analysis result.

    Two shapes count, because the Report pane sent notes as `{title: '', data: {conclusion}}` before the
    `note` field existed and a project saved then still opens.
    """
    return bool(section.note) or (not section.title and not (section.data or {}).get("tables"))


def _engine_meta(meta: ReportMeta, result_count: int = 0) -> report_builder.ReportMeta:
    """This module's ReportMeta as the engine's.

    Word and PowerPoint go through here too, purely for `meta_line`: the cover's mono metadata line is
    worded once, in the engine, so the four formats cannot disagree about it.
    """
    engine_meta = report_builder.ReportMeta(
        title=meta.title,
        dataset_label=meta.dataset_label,
        row_count=meta.row_count,
        column_count=meta.column_count,
        decimals=meta.decimals,
        version=APP_VERSION,
    )
    engine_meta.result_count = result_count
    return engine_meta


def _table_headers(rows: list[dict[str, Any]]) -> list[str]:
    headers: list[str] = []
    for row in rows:
        for key in row:
            if key not in headers:
                headers.append(key)
    return headers


def _numeric_columns(headers: list[str], rows: list[dict[str, Any]]) -> set[int]:
    """A column counts as numeric (right-aligned) when every filled value in it is a number."""
    numeric = set()
    for i, header in enumerate(headers):
        values = [row.get(header) for row in rows if row.get(header) is not None]
        if values and all(_is_number(v) for v in values):
            numeric.add(i)
    return numeric


def _write_captured_chart(data_url: str, stem: str, index: int, title: str) -> Path | None:
    """Decode a captured PNG to a file the renderers can embed. Bad input is skipped rather than
    failing the export — a missing figure is much better than a lost report."""
    try:
        payload = data_url.split(",", 1)[1] if "," in data_url else data_url
        raw = base64.b64decode(payload, validate=True)
        if not raw.startswith(bytes.fromhex("89504e47")):  # PNG magic number
            return None
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        dest = OUTPUT_DIR / f"{stem}_{index}_{_slugify(title)}_capture.png"
        # Flatten any alpha onto white: reportlab renders an RGBA PNG as an invisible rectangle,
        # so a transparent capture would silently vanish from the PDF.
        with PILImage.open(io.BytesIO(raw)) as image:
            if image.mode in ("RGBA", "LA", "P"):
                rgba = image.convert("RGBA")
                flat = PILImage.new("RGB", rgba.size, (255, 255, 255))
                flat.paste(rgba, mask=rgba.split()[-1])
                flat.save(dest, format="PNG")
            else:
                dest.write_bytes(raw)
        return dest
    except Exception:  # noqa: BLE001 - malformed capture, keep the rest of the report
        return None


def section_chart(section: ReportSection, stem: str, index: int) -> Path | None:
    """The PNG to embed under this section's table, in order of fidelity: the chart captured from
    the user's screen, then one generate_chart wrote to disk, then one rendered on demand from the
    analysis result."""
    if section.chart_image_base64:
        captured = _write_captured_chart(section.chart_image_base64, stem, index, section.title)
        if captured is not None:
            return captured
    if section.chart_path and Path(section.chart_path).exists():
        return Path(section.chart_path)
    if section.generated_chart is not None:
        return section.generated_chart if section.generated_chart.exists() else None
    if not section.allow_generated_chart:
        return None
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    dest = OUTPUT_DIR / f"{stem}_{index}_{_slugify(section.title)}_chart.png"
    try:
        produced = charts_core.render_analysis_chart(section.data, section.title, dest)
    except Exception:  # noqa: BLE001 - a chart is a nice-to-have; never fail the whole export
        produced = None
    section.generated_chart = produced
    return produced


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------


def _scalars_to_md(scalars: dict[str, Any], decimals: int) -> str:
    if not scalars:
        return ""
    lines = ["| Field | Value |", "|---|---|"]
    lines += [f"| {k} | {cell_text(v, decimals)} |" for k, v in scalars.items()]
    return "\n".join(lines)


def _rows_to_md(headers: list[str], rows: list[dict[str, Any]], numeric: set[int], decimals: int) -> str:
    if not rows:
        return ""
    align = ["---:" if i in numeric else "---" for i in range(len(headers))]
    lines = ["| " + " | ".join(headers) + " |", "|" + "|".join(align) + "|"]
    lines += ["| " + " | ".join(cell_text(row.get(h), decimals) for h in headers) + " |" for row in rows]
    return "\n".join(lines)


def render_markdown(meta: ReportMeta, sections: list[ReportSection], stem: str) -> Path:
    lines = [f"# {meta.title}", "", f"_{meta.shape_line}_", ""]

    for i, section in enumerate(sections, start=1):
        # A note carries no title: it is commentary in the flow of the report, so it gets prose and
        # no numbered heading.
        if section.title:
            lines += [f"## {i}. {section.title}", ""]
        _key, rows, scalars = _extract_table(section.data)
        narrative, rest = _split_narrative(scalars, section.title or "")

        # The PDF's verdict badge, in words. Same three fields, same wording, bolded rather than
        # coloured — so a reader skimming the Markdown gets the headline in the same place.
        badge = verdict_engine.line(verdict_engine.derive(section.data, section.analysis_id or section.title))
        if badge:
            lines += [f"**{badge}**", ""]

        if narrative:
            lines += [narrative, ""]
        if rows:
            shown, cut = _document_rows(rows, allow_full=section.full_tables)
            headers = _table_headers(shown)
            lines += [_rows_to_md(headers, shown, _numeric_columns(headers, shown), meta.decimals), ""]
            # Backticks: the note is mono in every Markdown viewer, and cannot be mistaken for a row.
            if cut:
                lines += [f"`{cut}`", ""]
        for table_title, table_rows in _extract_named_tables(section.data):
            shown, cut = _document_rows(table_rows, allow_full=section.full_tables)
            headers = _table_headers(shown)
            if table_title:
                lines += [f"**{table_title}**", ""]
            lines += [_rows_to_md(headers, shown, _numeric_columns(headers, shown), meta.decimals), ""]
            if cut:
                lines += [f"`{cut}`", ""]
        scalar_md = _scalars_to_md(rest, meta.decimals)
        if scalar_md:
            lines += [scalar_md, ""]

        chart = section_chart(section, stem, i)
        if chart:
            dest_name = f"{stem}_{i}_{_slugify(section.title)}.png"
            if chart.resolve() != (OUTPUT_DIR / dest_name).resolve():
                shutil.copyfile(chart, OUTPUT_DIR / dest_name)
            lines += [f"![{section.title}]({dest_name})", ""]

    lines += ["---", "", f"_{CREDIT_LINE}_", ""]
    md_path = OUTPUT_DIR / f"{stem}.md"
    md_path.write_text("\n".join(lines), encoding="utf-8")
    return md_path


# ---------------------------------------------------------------------------
# xlsx
# ---------------------------------------------------------------------------

_XL_THIN = Side(style="thin", color="FFD5DAE1")
_XL_BORDER = Border(left=_XL_THIN, right=_XL_THIN, top=_XL_THIN, bottom=_XL_THIN)


def _xlsx_safe(value: Any) -> Any:
    if value is None:
        return ""
    if isinstance(value, (list, dict)):
        return json.dumps(value)
    return value


def _xlsx_table(ws, start_row: int, rows: list[dict[str, Any]], decimals: int, widths: dict[int, int]) -> int:
    """Write one table starting at `start_row`; returns the first free row after it."""
    headers = _table_headers(rows)
    numeric = _numeric_columns(headers, rows)
    number_format = "0." + "0" * decimals if decimals > 0 else "0"
    row_cursor = start_row

    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=row_cursor, column=col_idx, value=header)
        cell.font = Font(bold=True, size=10, color="FF161616")
        cell.fill = PatternFill("solid", fgColor="FFF0F3F8")
        cell.border = _XL_BORDER
        cell.alignment = Alignment(horizontal="right" if (col_idx - 1) in numeric else "left")
        widths[col_idx] = max(widths.get(col_idx, 8), min(len(str(header)) + 3, 34))
    row_cursor += 1

    for row in rows:
        for col_idx, header in enumerate(headers, start=1):
            value = row.get(header)
            cell = ws.cell(row=row_cursor, column=col_idx, value=_xlsx_safe(value))
            cell.border = _XL_BORDER
            if _is_number(value):
                cell.alignment = Alignment(horizontal="right")
                if not float(value).is_integer():
                    # a p-value of 1e-12 under '0.000' would read as 0.000 in Excel too
                    cell.number_format = "0.00E+00" if value != 0 and abs(value) < 10 ** -decimals else number_format
            widths[col_idx] = max(widths.get(col_idx, 8), min(len(cell_text(value, decimals)) + 3, 34))
        row_cursor += 1
    return row_cursor + 1


def render_xlsx(meta: ReportMeta, sections: list[ReportSection], stem: str) -> Path:
    wb = Workbook()
    wb.remove(wb.active)
    used_names: set[str] = set()
    number_format = "0." + "0" * meta.decimals if meta.decimals > 0 else "0"

    for i, section in enumerate(sections, start=1):
        sheet_name = _sheet_name(section.title, i, used_names)
        used_names.add(sheet_name)
        ws = wb.create_sheet(title=sheet_name)

        _key, rows, scalars = _extract_table(section.data)
        narrative, rest = _split_narrative(scalars, section.title or "")

        title_cell = ws.cell(row=1, column=1, value=section.title)
        title_cell.font = Font(bold=True, size=13, color="FF161616")
        ws.cell(row=2, column=1, value=meta.shape_line).font = Font(size=9, color="FF6F6F6F")
        # A spreadsheet has no real page header, so the credit takes the third row of each sheet.
        ws.cell(row=3, column=1, value=CREDIT_LINE).font = Font(size=8, color="FF6F6F6F")
        row_cursor = 5

        # The verdict, in the same place the PDF puts its badge: bold, above the narrative, coloured
        # to the polarity. A cell's font colour is the closest a spreadsheet gets to a badge.
        fields = verdict_engine.derive(section.data, section.analysis_id or section.title)
        if fields:
            polarity = fields.get("verdict_polarity")
            colour = {"positive": _SUCCESS, "negative": _DANGER}.get(polarity, _MUTED)
            cell = ws.cell(row=row_cursor, column=1, value=verdict_engine.line(fields))
            cell.font = Font(bold=True, size=10.5, color=f"FF{colour}")
            row_cursor += 2

        if narrative:
            cell = ws.cell(row=row_cursor, column=1, value=narrative)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            ws.merge_cells(start_row=row_cursor, start_column=1, end_row=row_cursor + 2, end_column=8)
            row_cursor += 4

        widths: dict[int, int] = {}
        if rows:
            ws.freeze_panes = ws.cell(row=row_cursor + 1, column=1)
            row_cursor = _xlsx_table(ws, row_cursor, rows, meta.decimals, widths)

        for table_title, table_rows in _extract_named_tables(section.data):
            if table_title:
                label = ws.cell(row=row_cursor, column=1, value=table_title)
                label.font = Font(bold=True, size=10.5, color="FF161616")
                row_cursor += 1
            row_cursor = _xlsx_table(ws, row_cursor, table_rows, meta.decimals, widths)

        if widths:
            for col_idx, width in widths.items():
                ws.column_dimensions[get_column_letter(col_idx)].width = width

        if rest:
            ws.cell(row=row_cursor, column=1, value="Field").font = Font(bold=True, size=10)
            ws.cell(row=row_cursor, column=1).fill = PatternFill("solid", fgColor="FFF0F3F8")
            ws.cell(row=row_cursor, column=2, value="Value").font = Font(bold=True, size=10)
            ws.cell(row=row_cursor, column=2).fill = PatternFill("solid", fgColor="FFF0F3F8")
            row_cursor += 1
            for k, v in rest.items():
                ws.cell(row=row_cursor, column=1, value=k).border = _XL_BORDER
                cell = ws.cell(row=row_cursor, column=2, value=_xlsx_safe(v))
                cell.border = _XL_BORDER
                if _is_number(v):
                    cell.alignment = Alignment(horizontal="right")
                    if not float(v).is_integer():
                        cell.number_format = "0.00E+00" if v != 0 and abs(v) < 10 ** -meta.decimals else number_format
                row_cursor += 1
            ws.column_dimensions["A"].width = max(ws.column_dimensions["A"].width or 8, 24)
            row_cursor += 1

        chart = section_chart(section, stem, i)
        if chart:
            ws.add_image(XLImage(str(chart)), f"A{row_cursor + 1}")

    xlsx_path = OUTPUT_DIR / f"{stem}.xlsx"
    wb.save(xlsx_path)
    return xlsx_path


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------

_DOCX_PORTRAIT_WIDTH_IN = 6.5
_DOCX_LANDSCAPE_WIDTH_IN = 9.0
# Rough character budget for a table at 8.5pt: past this the table needs the wider orientation
# rather than columns squeezed to nothing.
_DOCX_PORTRAIT_CHAR_BUDGET = 88


# The same faces the PDF embeds, named as Word knows them. Word cannot be handed a TTF by python-docx
# — embedding needs the font obfuscated into an .odttf part — so these are NAMES: on a machine with the
# fonts installed the Word file is typographically identical to the PDF, and on one without it, Word
# substitutes and the structure (serif headings, mono metadata, the badge, the cards) still carries the
# hierarchy. Structure parity over pixel parity, deliberately.
_DOCX_SERIF = "Source Serif 4"
_DOCX_SERIF_FALLBACK = "Georgia"  # w:hAnsi, so a machine without Source Serif still gets A serif
_DOCX_SANS = "IBM Plex Sans"
_DOCX_MONO = "IBM Plex Mono"
_DOCX_MONO_FALLBACK = "Consolas"


def _docx_shade(cell, hex_fill: str) -> None:
    """python-docx has no cell-shading API, so set w:shd on the cell properties directly."""
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_fill.lstrip("#"))
    cell._tc.get_or_add_tcPr().append(shd)


def _docx_typeface(run, family: str, fallback: str = "") -> None:
    """Set the run's font on every script slot Word consults.

    `run.font.name` writes only w:ascii. Word picks w:hAnsi for anything above U+007F — which includes
    the badge's ● and the × in "96 rows × 8 columns" — so a run set the short way silently changes
    typeface mid-line at exactly those characters. Both slots, always.
    """
    rpr = run._r.get_or_add_rPr()
    fonts = rpr.find(qn("w:rFonts"))
    if fonts is None:
        fonts = OxmlElement("w:rFonts")
        rpr.insert(0, fonts)
    fonts.set(qn("w:ascii"), family)
    fonts.set(qn("w:hAnsi"), fallback or family)
    fonts.set(qn("w:cs"), family)


def _docx_run(paragraph, text: str, *, family: str = _DOCX_SERIF, fallback: str = _DOCX_SERIF_FALLBACK,
              size: float = 10, color: str = "", bold: bool = False, italic: bool = False):
    run = paragraph.add_run(text)
    _docx_typeface(run, family, fallback)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    return run


def _docx_border(element, edge: str, color: str, eighths: int = 6, space: int = 0) -> None:
    """One border edge on a paragraph or a table, in Word's own units (w:sz is EIGHTHS of a point)."""
    tag = "w:pBdr" if edge.startswith("p:") else "w:tblBorders"
    side = edge.split(":")[-1]
    container = element.find(qn(tag))
    if container is None:
        container = OxmlElement(tag)
        element.append(container)
    line = OxmlElement(f"w:{side}")
    line.set(qn("w:val"), "single")
    line.set(qn("w:sz"), str(eighths))
    line.set(qn("w:space"), str(space))
    line.set(qn("w:color"), color.lstrip("#"))
    container.append(line)


def _docx_card(container, width_in: float):
    """A bordered single-cell table standing in for the PDF's ResultCard, returning the cell.

    Word has no other way to draw a box around mixed content, and a table splits across a page break by
    itself with the border continuing — which is exactly the behaviour the PDF card was built for.
    """
    table = container.add_table(rows=1, cols=1)
    table.autofit = False
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    tbl_pr = table._tbl.tblPr
    for side in ("top", "left", "bottom", "right"):
        _docx_border(tbl_pr, f"t:{side}", _BORDER, eighths=6)
    # Cell padding: the PDF's 14pt, in twentieths of a point.
    margins = OxmlElement("w:tblCellMar")
    for side in ("top", "left", "bottom", "right"):
        node = OxmlElement(f"w:{side}")
        node.set(qn("w:w"), str(int(14 * 20)))
        node.set(qn("w:type"), "dxa")
        margins.append(node)
    tbl_pr.append(margins)
    cell = table.rows[0].cells[0]
    cell.width = Inches(width_in)
    # The single starting paragraph is empty; the caller writes into it first rather than after it, so
    # a card does not open with a blank line.
    return cell


def _docx_first_paragraph(cell):
    """The card's empty opening paragraph, so the first thing written lands at the top of the box."""
    return cell.paragraphs[0]


def _docx_picture(container, path: Path, *, width_in: float | None = None, height_in: float | None = None,
                  align=WD_ALIGN_PARAGRAPH.LEFT):
    """Insert a picture into a Document OR a table cell — `add_picture` only exists on Document."""
    paragraph = container.add_paragraph()
    paragraph.alignment = align
    run = paragraph.add_run()
    run.add_picture(str(path), width=Inches(width_in) if width_in else None,
                    height=Inches(height_in) if height_in else None)
    return paragraph


def _docx_repeat_header(row) -> None:
    tr_pr = row._tr.get_or_add_trPr()
    tbl_header = OxmlElement("w:tblHeader")
    tbl_header.set(qn("w:val"), "true")
    tr_pr.append(tbl_header)


def _docx_page_number(paragraph) -> None:
    """A live PAGE field, so the footer numbers pages instead of printing a fixed '1'."""
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.append(begin)
    run._r.append(instr)
    run._r.append(end)


def _docx_field(paragraph, instruction: str) -> None:
    """A live Word field (PAGE, NUMPAGES), so the footer counts pages instead of printing a fixed one."""
    run = paragraph.add_run()
    _docx_typeface(run, _DOCX_MONO, _DOCX_MONO_FALLBACK)
    run.font.size = Pt(7)
    run.font.color.rgb = RGBColor.from_string(_MUTED.lstrip("#"))
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = instruction
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.append(begin)
    run._r.append(instr)
    run._r.append(end)


def _docx_page_number(paragraph) -> None:
    """Kept as the narrow "just the PAGE field" helper the pptx/legacy callers use."""
    _docx_field(paragraph, "PAGE")


def _docx_header(section, title: str, assets: dict[str, Path], width_in: float) -> None:
    """The PDF's header, in Word: mark + wordmark left, title right, then the t-curve rule beneath.

    Word repeats a section header on every page of that section, which is how the signature ends up on
    every page here just as the reportlab onPage callback puts it on every page there.
    """
    section.header.is_linked_to_previous = False
    line = section.header.paragraphs[0]
    line.text = ""
    # One centre-less tab stop at the right margin puts the title hard right without a table.
    line.paragraph_format.tab_stops.add_tab_stop(Inches(width_in), WD_TAB_ALIGNMENT.RIGHT)
    if assets.get("mark"):
        line.add_run().add_picture(str(assets["mark"]), height=Pt(11))
        _docx_run(line, " ", size=11)
    _docx_run(line, APP_NAME, size=11.5, bold=True, color=_INK)
    _docx_run(line, "\t", size=7.5)
    _docx_run(line, title, family=_DOCX_MONO, fallback=_DOCX_MONO_FALLBACK, size=7.5, color=_MUTED)

    if assets.get("curve"):
        rule = section.header.add_paragraph()
        rule.paragraph_format.space_before = Pt(1)
        rule.paragraph_format.space_after = Pt(0)
        rule.add_run().add_picture(str(assets["curve"]), width=Inches(width_in))


def _docx_footer(section, assets: dict[str, Path], width_in: float) -> None:
    """The inverted curve, then the credit line and "page X of Y" — mirroring the PDF's footer."""
    section.footer.is_linked_to_previous = False
    first = section.footer.paragraphs[0]
    first.text = ""
    if assets.get("curve_inverted"):
        first.paragraph_format.space_after = Pt(2)
        first.add_run().add_picture(str(assets["curve_inverted"]), width=Inches(width_in))
        line = section.footer.add_paragraph()
    else:
        line = first
    line.paragraph_format.tab_stops.add_tab_stop(Inches(width_in), WD_TAB_ALIGNMENT.RIGHT)
    _docx_run(line, f"{CREDIT_LINE} v{APP_VERSION}\t", family=_DOCX_MONO, fallback=_DOCX_MONO_FALLBACK,
              size=7, color=_MUTED)
    _docx_field(line, "PAGE")
    _docx_run(line, " of ", family=_DOCX_MONO, fallback=_DOCX_MONO_FALLBACK, size=7, color=_MUTED)
    _docx_field(line, "NUMPAGES")


def _docx_orient(section, landscape: bool) -> None:
    width, height = (11, 8.5) if landscape else (8.5, 11)
    section.orientation = WD_ORIENT.LANDSCAPE if landscape else WD_ORIENT.PORTRAIT
    section.page_width = Inches(width)
    section.page_height = Inches(height)
    section.left_margin = section.right_margin = Inches(1.0)
    section.top_margin = section.bottom_margin = Inches(0.9)


def _docx_table(doc, headers: list[str], rows: list[dict[str, Any]], decimals: int, usable_width_in: float, font_size: float = 8.5):
    """One StatTable: no vertical rules, a shaded header row, mono right-aligned numerals.

    `doc` may be a Document or a table cell — both have add_table, which is what lets a table sit inside
    a card.
    """
    numeric = _numeric_columns(headers, rows)
    # The same per-column decimal consistency the PDF's StatTable uses, from the same function — so a
    # P column reads 0.000 in both documents rather than 0.000 in one and 0 in the other.
    cells = report_components.format_cells(rows, headers, decimals, cell_text)
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    table.autofit = False
    # Horizontal rules only, matching the PDF: a grid of boxes reads as a cage at this density. Set
    # explicitly rather than via the "Table Grid" style, which draws all four edges of every cell.
    tbl_pr = table._tbl.tblPr
    for side in ("top", "bottom", "insideH"):
        _docx_border(tbl_pr, f"t:{side}", _BORDER, eighths=4)

    weights = []
    for i, header in enumerate(headers):
        longest = len(str(header))
        for index in range(len(rows)):
            longest = max(longest, len(cells[(index, header)]))
        weights.append(min(max(longest, 4), 30))
    total = sum(weights) or 1
    widths = [max(0.6, usable_width_in * w / total) for w in weights]
    scale = usable_width_in / sum(widths)
    widths = [w * scale for w in widths]

    header_cells = table.rows[0].cells
    for i, header in enumerate(headers):
        cell = header_cells[i]
        cell.text = ""
        paragraph = cell.paragraphs[0]
        paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT if i in numeric else WD_ALIGN_PARAGRAPH.LEFT
        _docx_run(paragraph, str(header), family=_DOCX_SANS, fallback="Segoe UI", size=font_size,
                  bold=True, color=_INK)
        _docx_shade(cell, _HEADER_BG)
        cell.width = Inches(widths[i])
    _docx_repeat_header(table.rows[0])

    for row_index, _row in enumerate(rows):
        row_cells = table.add_row().cells
        for i, header in enumerate(headers):
            cell = row_cells[i]
            cell.text = ""
            paragraph = cell.paragraphs[0]
            paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT if i in numeric else WD_ALIGN_PARAGRAPH.LEFT
            # Mono for a number so the columns line up digit under digit, sans for a label: the same
            # split the PDF's StatTable makes.
            text = cells[(row_index, header)]
            if i in numeric:
                _docx_run(paragraph, text, family=_DOCX_MONO, fallback=_DOCX_MONO_FALLBACK,
                          size=font_size, color=_INK)
            else:
                _docx_run(paragraph, text, family=_DOCX_SANS, fallback="Segoe UI", size=font_size, color=_INK)
            cell.width = Inches(widths[i])
            if row_index % 2 == 1:
                _docx_shade(cell, _ROW_ALT_BG)
    return table


def _table_char_width(headers: list[str], rows: list[dict[str, Any]], decimals: int) -> int:
    total = 0
    for header in headers:
        longest = len(str(header))
        for row in rows:
            longest = max(longest, len(cell_text(row.get(header), decimals)))
        total += min(max(longest, 4), 30) + 2
    return total


def _docx_assets(stem: str, portrait_width_in: float) -> dict[str, Path]:
    """The mark and both curve rules, rasterised once per export from the report engine's own geometry.

    The curve is generated at the PORTRAIT text width and then stretched to the landscape width where
    needed: it is a smooth symmetric shape with no detail to lose, and one image beats two.
    """
    assets: dict[str, Path] = {}
    try:
        assets["mark"] = report_components.mark_png(OUTPUT_DIR / "_gosset_mark_ink.png", 128, f"#{_INK}")
    except Exception:  # noqa: BLE001 - the header still works without it; never lose a report over a logo
        pass
    try:
        # Peak left to the engine's own CURVE_PEAK, so retuning the rule reaches Word too.
        width_pt = portrait_width_in * 72
        assets["curve"] = report_components.curve_png(
            OUTPUT_DIR / "_gosset_curve.png", width_pt, color=f"#{_INK}")
        assets["curve_inverted"] = report_components.curve_png(
            OUTPUT_DIR / "_gosset_curve_inv.png", width_pt, invert=True, color=f"#{_INK}")
    except Exception:  # noqa: BLE001 - same: the rule is the signature, not the content
        pass
    return assets


def render_docx(meta: ReportMeta, sections: list[ReportSection], stem: str) -> Path:
    """The Word mirror of the PDF: same hierarchy, same words, within python-docx's limits.

    What carries over: the every-page header (mark, wordmark, t-curve rule), the inverted footer rule
    with page X of Y, a serif cover title over a mono metadata line, one bordered card per result, the
    verdict badge, mono numerals in the tables, captions under the charts and notes as indented prose.
    What does not: rounded card corners and true small caps, neither of which Word offers.
    """
    doc = Document()
    assets = _docx_assets(stem, _DOCX_PORTRAIT_WIDTH_IN)
    base = doc.sections[0]
    _docx_orient(base, landscape=False)
    _docx_header(base, meta.title, assets, _DOCX_PORTRAIT_WIDTH_IN)
    _docx_footer(base, assets, _DOCX_PORTRAIT_WIDTH_IN)

    # Body text is the serif, matching the PDF, and every heading style follows it — so a reader who
    # restyles the document keeps the hierarchy.
    normal = doc.styles["Normal"]
    normal.font.name = _DOCX_SERIF
    normal.font.size = Pt(10)
    for style_name in ("Title", "Heading 1", "Heading 2", "Heading 3"):
        try:
            style = doc.styles[style_name]
        except KeyError:
            continue
        style.font.name = _DOCX_SERIF
        style.font.color.rgb = RGBColor.from_string(_INK.lstrip("#"))

    # --- the cover block, on the first page above the first card ---
    title_para = doc.add_paragraph()
    title_para.paragraph_format.space_after = Pt(2)
    _docx_run(title_para, meta.title, size=23, bold=True, color=_INK)
    meta_para = doc.add_paragraph()
    meta_para.paragraph_format.space_after = Pt(18)
    described = len([s for s in sections if not _is_note_section(s)])
    _docx_run(meta_para, _engine_meta(meta, described).meta_line, family=_DOCX_MONO,
              fallback=_DOCX_MONO_FALLBACK, size=8, color=_MUTED)

    landscape_now = False
    for i, section in enumerate(sections, start=1):
        # --- a Report-pane note: the author speaking, not a result. Prose, indented, no card ---
        if _is_note_section(section):
            text = section.note or str((section.data or {}).get("conclusion") or "")
            if not text.strip():
                continue
            note = doc.add_paragraph()
            note.paragraph_format.left_indent = Inches(0.28)
            note.paragraph_format.space_before = Pt(6)
            note.paragraph_format.space_after = Pt(10)
            _docx_border(note._p.get_or_add_pPr(), "p:left", _ACCENT, eighths=12, space=10)
            _docx_run(note, text.strip(), size=10, color=_INK)
            continue

        _key, rows, scalars = _extract_table(section.data)
        narrative, rest = _split_narrative(scalars, section.title or "")
        headers = _table_headers(rows) if rows else []
        named = _extract_named_tables(section.data)

        # Decide the orientation before writing anything, so the heading and its narrative travel
        # into the same section as the table instead of being stranded on the previous page. The
        # widest table in the section decides, not just the first.
        widths_needed = [_table_char_width(headers, rows, meta.decimals)] if rows else []
        widths_needed += [_table_char_width(_table_headers(r), r, meta.decimals) for _t, r in named]
        wants_landscape = bool(widths_needed) and max(widths_needed) > _DOCX_PORTRAIT_CHAR_BUDGET
        if wants_landscape != landscape_now:
            new_section = doc.add_section(WD_SECTION.NEW_PAGE)
            _docx_orient(new_section, landscape=wants_landscape)
            page_width = _DOCX_LANDSCAPE_WIDTH_IN if wants_landscape else _DOCX_PORTRAIT_WIDTH_IN
            _docx_header(new_section, meta.title, assets, page_width)
            _docx_footer(new_section, assets, page_width)
            landscape_now = wants_landscape
        current_width = _DOCX_LANDSCAPE_WIDTH_IN if landscape_now else _DOCX_PORTRAIT_WIDTH_IN

        # --- the card ---
        card = _docx_card(doc, current_width)
        inner_width = current_width - 0.42  # the 14pt padding on both sides
        heading = _docx_first_paragraph(card)
        heading.style = doc.styles["Heading 1"]  # inside a cell, and Word's navigation pane still sees it
        heading.paragraph_format.space_before = Pt(0)
        heading.paragraph_format.space_after = Pt(1)
        _docx_run(heading, section.title, size=13.5, bold=True, color=_INK)

        card_meta = "  ·  ".join(b for b in (section.timestamp, section.columns) if b)
        if card_meta:
            line = card.add_paragraph()
            line.paragraph_format.space_after = Pt(4)
            _docx_run(line, card_meta, family=_DOCX_MONO, fallback=_DOCX_MONO_FALLBACK, size=7.5, color=_MUTED)
        # The hairline under the card header, as a paragraph's bottom border.
        rule = card.add_paragraph()
        rule.paragraph_format.space_before = Pt(0)
        rule.paragraph_format.space_after = Pt(6)
        _docx_border(rule._p.get_or_add_pPr(), "p:bottom", _HAIRLINE, eighths=4, space=1)

        fields = verdict_engine.derive(section.data, section.analysis_id or section.title)
        if fields:
            badge = card.add_paragraph()
            badge.paragraph_format.space_after = Pt(6)
            polarity = fields.get("verdict_polarity")
            colour = {"positive": _SUCCESS, "negative": _DANGER}.get(polarity, _MUTED)
            # ● in the polarity colour, then the label bold, then the statistic in mono. Word's fonts do
            # carry U+25CF, so unlike the PDF this really is the glyph rather than a drawn circle.
            _docx_run(badge, "● ", family=_DOCX_SANS, fallback="Segoe UI Symbol", size=10, color=colour)
            _docx_run(badge, fields.get("verdict_label", ""), family=_DOCX_SANS, fallback="Segoe UI",
                      size=10, bold=True, color=_INK)
            if fields.get("key_stat"):
                _docx_run(badge, f"   {fields['key_stat']}", family=_DOCX_MONO,
                          fallback=_DOCX_MONO_FALLBACK, size=9.5, color=_INK_80)

        def _docx_table_note(text: str) -> None:
            """The mono "… N more rows" line under a table that was cut short."""
            if not text:
                return
            line = card.add_paragraph()
            line.paragraph_format.space_before = Pt(1)
            line.paragraph_format.space_after = Pt(2)
            _docx_run(line, text, family=_DOCX_MONO, fallback=_DOCX_MONO_FALLBACK, size=7.5, color=_MUTED)

        if rows:
            shown, cut = _document_rows(rows, allow_full=section.full_tables)
            _docx_table(card, _table_headers(shown), shown, meta.decimals, inner_width)
            _docx_table_note(cut)
            card.add_paragraph().paragraph_format.space_after = Pt(2)

        for table_title, table_rows in named:
            if table_title:
                label = card.add_paragraph()
                label.paragraph_format.space_after = Pt(2)
                # The PDF's small-caps group label. Word has no small-caps-from-lowercase without a real
                # small-caps face, so it is upper-cased and letterspaced by eye instead.
                _docx_run(label, table_title.upper(), family=_DOCX_MONO, fallback=_DOCX_MONO_FALLBACK,
                          size=7, bold=True, color=_MUTED)
            shown, cut = _document_rows(table_rows, allow_full=section.full_tables)
            _docx_table(card, _table_headers(shown), shown, meta.decimals, inner_width)
            _docx_table_note(cut)
            card.add_paragraph().paragraph_format.space_after = Pt(2)

        if rest:
            _docx_table(card, ["Field", "Value"], [{"Field": k, "Value": v} for k, v in rest.items()],
                        meta.decimals, min(inner_width, 5.2))
            card.add_paragraph().paragraph_format.space_after = Pt(2)

        chart = section_chart(section, stem, i)
        if chart:
            # Word does not shrink an oversized picture, so a tall stacked capture has to be sized
            # by its height instead of its width or it runs off the page.
            picture_width = min(inner_width, 6.1)
            max_height = 5.6 if landscape_now else 7.8
            try:
                with PILImage.open(chart) as image:
                    tall = image.height / image.width * picture_width > max_height
            except Exception:  # noqa: BLE001 - unreadable image, let python-docx decide
                tall = False
            if tall:
                _docx_picture(card, chart, height_in=max_height, align=WD_ALIGN_PARAGRAPH.CENTER)
            else:
                _docx_picture(card, chart, width_in=picture_width, align=WD_ALIGN_PARAGRAPH.CENTER)

        if narrative:
            cap = card.add_paragraph()
            cap.paragraph_format.space_before = Pt(4)
            cap.paragraph_format.space_after = Pt(0)
            _docx_run(cap, narrative, family=_DOCX_SANS, fallback="Segoe UI", size=8.5, color=_INK_80)

        # Air between cards. Outside the card, so it is a gap and not padding.
        doc.add_paragraph().paragraph_format.space_after = Pt(4)

    docx_path = OUTPUT_DIR / f"{stem}.docx"
    doc.save(docx_path)
    return docx_path


# ---------------------------------------------------------------------------
# pdf
# ---------------------------------------------------------------------------

def _section_tables(data: dict[str, Any], *, allow_full: bool = False) -> list[tuple[str, list[dict[str, Any]], str]]:
    """Every table in one result as [(group label, rows, truncation note)], in reading order.

    A result may ship its tables three different ways and the engine should not have to know that:
    the `tables: [{title, rows}]` list every procedure returns, a bare priority list like
    `coefficients`, and the leftover scalars. All three collapse to the same triples here.

    Long tables are cut to TABLE_ROWS_SHOWN with a note saying how many rows were left out, unless the
    section was staged as a full table. The note travels as the third element so the engine can set it
    in mono under the table it belongs to; older two-element pairs still unpack in the builder.
    """
    tables: list[tuple[str, list[dict[str, Any]], str]] = []
    named = _extract_named_tables(data)
    _key, rows, scalars = _extract_table(data)
    _narrative, rest = _split_narrative(scalars)

    if named:
        for label, table_rows in named:
            shown, note = _document_rows(table_rows, allow_full=allow_full)
            tables.append((label, shown, note))
    elif rows:
        shown, note = _document_rows(rows, allow_full=allow_full)
        tables.append(("", shown, note))
    if rest:
        # The loose scalars as a two-column table. Titled so it reads as a group rather than as an
        # unexplained pair of columns after the real results. Never truncated: these are the run's own
        # settings and there are never many.
        tables.append(("Details", [{"Field": k, "Value": v} for k, v in rest.items()], ""))
    return tables


def _engine_sections(meta: ReportMeta, sections: list[ReportSection], stem: str) -> list["report_builder.Section"]:
    """Translate this module's ReportSection into the report engine's Section.

    The split of responsibility: everything that knows the SHAPE of an analysis result — which key
    holds the table, which fields are narrative — stays here, because that is app knowledge. The
    engine is handed titles, rows, a PNG path and a caption, and knows nothing about statistics.
    """
    built: list[report_builder.Section] = []
    for index, section in enumerate(sections, start=1):
        if _is_note_section(section):
            # A Report-pane note. `data.conclusion` is where the pane used to put its text, so both
            # spellings are accepted and a project saved before the `note` field existed still reads.
            text = section.note or str((section.data or {}).get("conclusion") or "")
            if text.strip():
                built.append(report_builder.Section(note=text.strip()))
            continue

        _key, _rows, scalars = _extract_table(section.data)
        narrative, _rest = _split_narrative(scalars, section.title or "")
        chart = section_chart(section, stem, index)
        built.append(
            report_builder.Section(
                title=section.title,
                # The title is a good fallback subject for the badge's polarity rule: "Normality Test —
                # residuals" normalises to a string containing `normality` just as the id does.
                analysis_id=section.analysis_id or section.title,
                data=section.data or {},
                tables=_section_tables(section.data, allow_full=section.full_tables),
                chart=chart,
                caption=narrative,
                columns=section.columns,
                timestamp=section.timestamp,
            )
        )
    return built


def render_pdf(meta: ReportMeta, sections: list[ReportSection], stem: str) -> Path:
    """The branded PDF, via backend/report_engine.

    This function is now a translator and nothing else: the layout, the fonts and the t-curve
    furniture all live in the engine, which has no idea what a p-value is.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pdf_path = OUTPUT_DIR / f"{stem}.pdf"
    return report_builder.render_pdf(
        pdf_path, _engine_sections(meta, sections, stem), _engine_meta(meta), cell_text)


# ---------------------------------------------------------------------------
# entry point
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------


def _pptx_table(slide, rows: list[dict[str, Any]], decimals: int, top: float, height: float, width: float, left: float):
    """One table, styled to match the app's tables: accent header, hairline rules, numbers right."""
    headers = _table_headers(rows)
    numeric = _numeric_columns(headers, rows)
    # A slide is not a page: past ~12 rows the type shrinks to nothing, so the table is truncated and
    # says so rather than rendering an illegible wall.
    shown = rows[:12]
    shape = slide.shapes.add_table(len(shown) + 1, len(headers), left, top, width, height)
    table = shape.table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = str(header)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = PptPt(11)
        para.font.color.rgb = PptRGB.from_string(_INK.lstrip("#"))
        para.alignment = PP_ALIGN.RIGHT if col in numeric else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGB.from_string(_HEADER_BG.lstrip("#"))
    for r, row in enumerate(shown, start=1):
        for col, header in enumerate(headers):
            cell = table.cell(r, col)
            cell.text = cell_text(row.get(header), decimals)
            para = cell.text_frame.paragraphs[0]
            para.font.size = PptPt(10)
            para.font.color.rgb = PptRGB.from_string(_INK.lstrip("#"))
            para.alignment = PP_ALIGN.RIGHT if col in numeric else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGB.from_string("FFFFFF")
    return len(rows) - len(shown)


def render_pptx(meta: ReportMeta, sections: list[ReportSection], stem: str) -> Path:
    """One slide per section: the title, then the chart if there is one, otherwise the table.

    Built for the block menu's "Send to PowerPoint", where a section is a single output block, so a
    slide holds one thing and holds it big. Passing several sections gives several slides, which is
    what the Report pane's export does.
    """
    prs = Presentation()
    prs.slide_width = PptInches(13.333)  # 16:9
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]  # no placeholders; everything is positioned explicitly

    margin = PptInches(0.6)
    content_w = prs.slide_width - 2 * margin

    for i, section in enumerate(sections, start=1):
        slide = prs.slides.add_slide(blank)

        title_box = slide.shapes.add_textbox(margin, PptInches(0.42), content_w, PptInches(0.7))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = section.title
        title_para.font.size = PptPt(26)
        title_para.font.bold = True
        title_para.font.color.rgb = PptRGB.from_string(_INK.lstrip("#"))

        sub_box = slide.shapes.add_textbox(margin, PptInches(1.02), content_w, PptInches(0.34))
        sub_para = sub_box.text_frame.paragraphs[0]
        sub_para.text = meta.dataset_label
        sub_para.font.size = PptPt(12)
        sub_para.font.color.rgb = PptRGB.from_string(_MUTED.lstrip("#"))

        _key, rows, scalars = _extract_table(section.data)
        # A block sent from the UI carries its one table under `tables: [{title, rows}]` rather than
        # under one of the priority keys, so _extract_table finds nothing. Without this fallback the
        # slide silently degrades to just the narrative — which is what the first version did.
        if not rows:
            named = _extract_named_tables(section.data)
            if named:
                rows = named[0][1]
        narrative, _rest = _split_narrative(scalars, section.title or "")
        body_top = PptInches(1.55)
        body_h = prs.slide_height - body_top - PptInches(0.95)

        chart = section_chart(section, stem, i)
        if chart is not None and chart.exists():
            # Fit the image inside the body box, centred, without distorting it.
            with PILImage.open(chart) as img:
                iw, ih = img.size
            scale = min(content_w / iw, body_h / ih)
            w, h = int(iw * scale), int(ih * scale)
            slide.shapes.add_picture(str(chart), margin + (content_w - w) // 2, body_top + (body_h - h) // 2, width=w, height=h)
        elif rows:
            table_h = min(body_h, PptInches(0.32) * (min(len(rows), 12) + 1))
            dropped = _pptx_table(slide, rows, meta.decimals, body_top, table_h, content_w, margin)
            if dropped > 0:
                note = slide.shapes.add_textbox(margin, body_top + table_h + PptInches(0.1), content_w, PptInches(0.3))
                note_para = note.text_frame.paragraphs[0]
                note_para.text = f"+ {dropped} more row(s) — see the full report for all of them."
                note_para.font.size = PptPt(10)
                note_para.font.color.rgb = PptRGB.from_string(_MUTED.lstrip("#"))
        elif narrative:
            box = slide.shapes.add_textbox(margin, body_top, content_w, body_h)
            frame = box.text_frame
            frame.word_wrap = True
            frame.paragraphs[0].text = narrative
            frame.paragraphs[0].font.size = PptPt(16)
            frame.paragraphs[0].font.color.rgb = PptRGB.from_string(_INK.lstrip("#"))

        if narrative and (chart is not None or rows):
            # The narrative becomes the takeaway line under the visual, which is what a slide wants.
            box = slide.shapes.add_textbox(margin, prs.slide_height - PptInches(0.92), content_w, PptInches(0.42))
            frame = box.text_frame
            frame.word_wrap = True
            frame.paragraphs[0].text = narrative
            frame.paragraphs[0].font.size = PptPt(12)
            frame.paragraphs[0].font.color.rgb = PptRGB.from_string(_INK.lstrip("#"))

        credit = slide.shapes.add_textbox(margin, prs.slide_height - PptInches(0.46), content_w, PptInches(0.3))
        credit_para = credit.text_frame.paragraphs[0]
        credit_para.text = CREDIT_LINE
        credit_para.font.size = PptPt(9)
        credit_para.font.color.rgb = PptRGB.from_string(_MUTED.lstrip("#"))

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    dest = OUTPUT_DIR / f"{stem}.pptx"
    prs.save(str(dest))
    return dest


RENDERERS = {
    "markdown": render_markdown,
    "xlsx": render_xlsx,
    "docx": render_docx,
    "pdf": render_pdf,
    "pptx": render_pptx,
}


def build_report(meta: ReportMeta | str, fmt: str, sections: list[ReportSection], stem: str) -> list[Path]:
    """`fmt` is one of markdown / xlsx / docx / pdf, or "both" (markdown + xlsx), kept because the
    MCP tool has shipped with that value."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    report_meta = coerce_meta(meta)
    wanted = ["markdown", "xlsx"] if fmt == "both" else [fmt]
    paths: list[Path] = []
    for name in wanted:
        renderer = RENDERERS.get(name)
        if renderer is None:
            raise ValueError(f"Unsupported report format '{fmt}'. Expected one of: {', '.join(RENDERERS)}, both.")
        paths.append(renderer(report_meta, sections, stem))
    return paths
